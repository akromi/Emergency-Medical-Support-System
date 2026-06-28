#!/usr/bin/env python3
"""Generate the TRIAGE-LINK capability deck (.pptx) in the app's dark theme."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- brand palette (from src/styles.css) ----
BG     = RGBColor(0x0E, 0x11, 0x16)
PANEL  = RGBColor(0x16, 0x1B, 0x22)
PANEL2 = RGBColor(0x1B, 0x22, 0x2C)
LINE   = RGBColor(0x2A, 0x33, 0x40)
INK    = RGBColor(0xE6, 0xED, 0xF3)
DIM    = RGBColor(0xA2, 0xAE, 0xBE)
FAINT  = RGBColor(0x8E, 0x9B, 0xAD)
EKG    = RGBColor(0x3F, 0xE0, 0x8A)   # signature green
AMBER  = RGBColor(0xF2, 0x82, 0x0F)
RED    = RGBColor(0xE5, 0x48, 0x4D)
BLUE   = RGBColor(0x3F, 0x9B, 0xD6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
SANS = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bg(slide, color=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    _solid(r, color)
    r.shadow.inherit = False
    return r


def rect(slide, x, y, w, h, color, line=None, radius=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    _solid(shp, color)
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (text,size,color,bold,font)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.line_spacing = line_spacing
        for (t, size, color, bold, font) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
    return tb


def footer(slide, n):
    text(slide, Inches(0.55), Inches(7.02), Inches(8), Inches(0.3),
         [[("TRIAGE-LINK", 9, EKG, True, MONO),
           ("   offline-first field casualty documentation", 9, FAINT, False, SANS)]])
    text(slide, Inches(11.6), Inches(7.02), Inches(1.2), Inches(0.3),
         [[(f"{n:02d}", 9, FAINT, False, MONO)]], align=PP_ALIGN.RIGHT)


N = [0]

def chrome(slide, kicker, title, accent=EKG):
    N[0] += 1
    bg(slide)
    rect(slide, Inches(0.55), Inches(0.6), Inches(0.09), Inches(0.92), accent)
    text(slide, Inches(0.82), Inches(0.55), Inches(11.8), Inches(0.35),
         [[(kicker.upper(), 12, accent, True, MONO)]])
    text(slide, Inches(0.8), Inches(0.86), Inches(11.9), Inches(0.8),
         [[(title, 30, INK, True, SANS)]])
    footer(slide, N[0])


def bullets(slide, items, x=Inches(0.85), y=Inches(1.95), w=Inches(11.7), h=Inches(4.9),
            size=16, gap=9):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        if lvl == 0:
            runs.append([("▸  ", 16, EKG, True, SANS), (txt, size, INK, False, SANS)])
        else:
            runs.append([("      –  ", 14, FAINT, False, SANS), (txt, size-2, DIM, False, SANS)])
    text(slide, x, y, w, h, runs, space=gap, line_spacing=1.05)


def two_col(slide, left_title, left_items, right_title, right_items, accentL=EKG, accentR=AMBER):
    cardw, cardh = Inches(5.75), Inches(4.7)
    for (cx, ctitle, citems, acc) in [
        (Inches(0.7), left_title, left_items, accentL),
        (Inches(6.9), right_title, right_items, accentR)]:
        rect(slide, cx, Inches(1.95), cardw, cardh, PANEL, line=LINE, radius=True)
        text(slide, cx+Inches(0.3), Inches(2.18), cardw-Inches(0.6), Inches(0.4),
             [[(ctitle, 15, acc, True, SANS)]])
        runs = [[("•  ", 13, acc, True, SANS), (t, 13.5, INK, False, SANS)] for t in citems]
        text(slide, cx+Inches(0.3), Inches(2.75), cardw-Inches(0.55), Inches(3.8),
             runs, space=7, line_spacing=1.04)


# ============================ SLIDES ============================

# 1 — Title
s = prs.slides.add_slide(BLANK)
N[0] += 1
bg(s)
rect(s, 0, Inches(3.42), EMU_W, Inches(0.045), EKG)
text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.2),
     [[("TRIAGE-LINK", 60, INK, True, SANS)]])
text(s, Inches(0.95), Inches(3.6), Inches(11.5), Inches(0.6),
     [[("Offline-first PWA for field casualty documentation & coordination", 22, EKG, False, SANS)]])
text(s, Inches(0.95), Inches(4.35), Inches(11.5), Inches(1.4),
     [[("Capability overview — the PWA, the encryption & audit layer, the hosted ", 15, DIM, False, SANS)],
      [("multi-tenant backend with admin security, and the market flavors ", 15, DIM, False, SANS)],
      [("(Humanitarian / NGO · Ontario EMS / regulated · productized backend).", 15, DIM, False, SANS)]], space=3)
text(s, Inches(0.95), Inches(6.7), Inches(11.5), Inches(0.4),
     [[("React + TypeScript · IndexedDB + op-log · WebCrypto · Fastify · NEMSIS/OADS", 12, FAINT, False, MONO)]])

# 2 — What it is
s = prs.slides.add_slide(BLANK)
chrome(s, "Positioning", "One field record, no signal required")
bullets(s, [
    "Document a casualty completely OFFLINE — injuries, vitals, treatments, photos, triage, handover — with no server in the loop.",
    "Built for where connectivity is unreliable or absent: disaster/MCI response, humanitarian field clinics, mass-gathering and prehospital EMS.",
    "Installs as a PWA on any device; data lives locally in IndexedDB and survives reloads, crashes, and power cycles.",
    "Optional, never-required sync aggregates records across a team only where a deployment wants it.",
    "One codebase, three market 'flavors' on branches that share the offline-first core.",
])

# 3 — Architecture
s = prs.slides.add_slide(BLANK)
chrome(s, "Architecture", "Offline-first core, optional everything else")
boxes = [
    (Inches(0.7), "DEVICE — PWA", EKG, [
        "React + TS UI", "IndexedDB (Dexie)", "Op-log sync engine",
        "WebCrypto vault", "Operator roster + audit"]),
    (Inches(5.0), "OPTIONAL — Sync service", AMBER, [
        "Fastify, multi-tenant", "Conflict-aware /sync", "OIDC admin + console",
        "Quota · retention · rate-limit", "Postgres"]),
    (Inches(9.3), "OPTIONAL — EHR gateway", BLUE, [
        "ONE ID / Ontario Health", "PCR $match · handover", "mTLS, server-side only",
        "NEMSIS/OADS export", "Audit trail"]),
]
for (x, title, acc, items) in boxes:
    rect(s, x, Inches(2.3), Inches(3.3), Inches(3.6), PANEL, line=LINE, radius=True)
    rect(s, x, Inches(2.3), Inches(3.3), Inches(0.12), acc, radius=False)
    text(s, x+Inches(0.25), Inches(2.5), Inches(2.9), Inches(0.5), [[(title, 13, acc, True, MONO)]])
    text(s, x+Inches(0.25), Inches(3.05), Inches(2.9), Inches(2.7),
         [[("•  ", 12, acc, True, SANS), (t, 12.5, INK, False, SANS)] for t in items],
         space=6, line_spacing=1.05)
for ax in [Inches(4.32), Inches(8.62)]:
    text(s, ax, Inches(3.7), Inches(0.66), Inches(0.5), [[("→", 26, FAINT, True, SANS)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.7),
     [[("The PWA is fully functional with neither optional tier present. Sync and the EHR gateway are additive, "
        "server-side, and off by default.", 13, DIM, False, SANS)]])

# 4 — PWA documenting
s = prs.slides.add_slide(BLANK)
chrome(s, "PWA · Capture", "Documenting a casualty")
two_col(s, "Body & injuries", [
    "Tap-to-mark injuries on an anterior/posterior body chart",
    "Injury type palette (GSW, burn, laceration, fracture…) with severity",
    "Per-injury notes + wound photos (stored out-of-line as blobs)",
    "Burn TBSA auto-estimate (Lund–Browder by age band)",
], "Triage & identity", [
    "START-style triage tag: Immediate / Delayed / Minor / Deceased",
    "Patient tombstone: name, DOB→age band, sex, MRN, NOK, blood type",
    "Incident: time, mechanism, location",
    "Color-coded multi-casualty Triage Board (the scene picture)",
])

# 5 — vitals/treatments
s = prs.slides.add_slide(BLANK)
chrome(s, "PWA · Clinical", "Vitals, interventions & trends")
two_col(s, "Vitals & scoring", [
    "Timestamped vital sets: HR, BP, RR, SpO₂, GCS, pain",
    "Built-in GCS calculator (eye/verbal/motor → total)",
    "Vitals-trend sparklines once ≥2 readings exist",
    "Time-since-injury clock (T+) on the record",
], "Treatments", [
    "Structured treatment log: tourniquet, airway, decompression,",
    "  IV/fluids, medication, splinting, wound packing, CPR…",
    "Each entry timestamped and attributed to the operator on duty",
    "Feeds the AT-MIST handover summary",
])

# 6 — handover/summaries
s = prs.slides.add_slide(BLANK)
chrome(s, "PWA · Handover", "Summaries & clean handoff")
bullets(s, [
    "One-page printable Casualty Summary card (AT-MIST) — print or save as PDF for handover.",
    "Scene Summary / command roll-up: casualties tallied by triage, on-scene vs handed-over.",
    "Handover sign-off (who took over care, facility, time) emitted as a FHIR handover bundle.",
    "Optional 'Send to EHR' contributes the handover to a provincial EHR via the gateway.",
    "Everything renders offline; nothing leaves the device unless you export or sync.",
])

# 7 — i18n + tour
s = prs.slides.add_slide(BLANK)
chrome(s, "PWA · Accessibility", "Four languages, guided & spoken")
two_col(s, "Internationalization", [
    "EN / FR / AR / FA built in — Arabic & Persian fully RTL",
    "Loadable JSON language packs: add a language with NO app release",
    "Downloadable English template to translate; parity-tested in CI",
    "Natural wording, not literal — reviewed per language",
], "Guided tour", [
    "Smart guided tour highlights each real control",
    "Offline voice-over (SpeechSynthesis) in the active language",
    "Action steps auto-advance once the user does them",
    "Every user-visible feature is taught in the tour (enforced)",
], accentR=BLUE)

# 8 — offline/install
s = prs.slides.add_slide(BLANK)
chrome(s, "PWA · Platform", "Installable, offline, durable")
bullets(s, [
    "Installable PWA (service worker + Workbox precache) — launches and runs with no network.",
    "All data in IndexedDB via Dexie; records, op-log, photos, audit chain persist locally.",
    "Works on phones, tablets, and laptops; responsive layout collapses gracefully on small screens.",
    "No telemetry, no implicit network calls — a casualty is documented entirely on-device.",
])

# 9 — op-log
s = prs.slides.add_slide(BLANK)
chrome(s, "Data integrity", "Conflict-aware op-log sync engine", accent=BLUE)
bullets(s, [
    "Every change is journaled as an immutable operation (scalar fields + collections), not a blind overwrite.",
    "Deterministic resolve(): Lamport clocks order edits; ties break predictably — same inputs, same result, on every device.",
    "Concurrent edits to different fields all survive; same-field edits pick a deterministic winner and REPORT the conflict (losing op retained, never silently dropped).",
    "The server stores and folds ops — it does not implement its own divergent merge logic.",
    "Incremental sync (cursor) pulls only what changed; full-state pulls are paginated.",
])

# 10 — vault
s = prs.slides.add_slide(BLANK)
chrome(s, "Security · At rest", "Opt-in encryption vault", accent=AMBER)
two_col(s, "How it works", [
    "AES-256-GCM, key derived from a passphrase via PBKDF2 (210k iters)",
    "Encrypts the heaviest PHI — wound photos — plus records & op-log",
    "Key lives only in memory while unlocked; locking drops it",
    "Idle auto-lock; wrong passphrase rejected via a verifier, no data touched",
], "Safety posture", [
    "DEFAULT-OFF: with no vault, behavior is byte-for-byte unchanged",
    "Mixed plaintext/encrypted rows read correctly through a toggle",
    "Sealed records are unreadable (get/list skip them) while locked",
    "Crash-safe enable/disable — data is never orphaned mid-migration",
], accentL=AMBER, accentR=EKG)

# 11 — operators/audit
s = prs.slides.add_slide(BLANK)
chrome(s, "Security · Access", "Operators, RBAC-lite & tamper-evident audit", accent=AMBER)
two_col(s, "Shared-device access", [
    "Local operator roster (field / lead / admin roles)",
    "Records & audit entries attributed to the on-duty operator",
    "Step-up PIN re-auth gates sensitive actions (delete, export…)",
    "Empty roster = open (community default); adding operators opts in",
], "Audit log", [
    "Append-only, hash-chained entries (SHA-256 over each entry)",
    "Tampering with or deleting any entry breaks the chain — detectable offline",
    "No update/delete API; reviewable even while the vault is locked",
    "Covers create / view / delete / export / vault / step-up events",
], accentL=AMBER, accentR=AMBER)

# 12 — backup/export
s = prs.slides.add_slide(BLANK)
chrome(s, "Data portability", "Backup, restore & export")
two_col(s, "Backup / restore", [
    "Full JSON backup of every record (plain or passphrase-encrypted)",
    "Restore by merge (keep newer of duplicates) or replace",
    "Encrypted backup keeps PHI unreadable without the passphrase",
], "CSV interchange", [
    "Roster CSV export (identity + incident fields) for analytics/QA",
    "CSV import onboards a patient list from paper or another system",
    "Date-range filter scopes an export to a time window",
    "Deployment provenance stamped on every row (humanitarian flavor)",
], accentR=BLUE)

# 13 — backend
s = prs.slides.add_slide(BLANK)
chrome(s, "Hosted backend", "Multi-tenant sync service (Fastify)", accent=AMBER)
bullets(s, [
    "Optional cloud or self-hosted backend for cross-team aggregation — the PWA never requires it.",
    "Per-tenant isolation: each API key authenticates AND scopes a tenant's data; conflict-aware /sync stores and folds ops.",
    "Hardened: sanitized error envelopes (no internal leakage), paginated full-state pulls, per-tenant storage quota (noisy-neighbor guard), audit-log retention TTL.",
    "OpenAPI 3 document + Swagger UI; liveness/readiness probes; per-tenant operational metrics; request-id correlation.",
    "Runs on PostgreSQL; graceful shutdown; rate-limited per IP.",
], size=15, gap=10)

# 14 — admin security
s = prs.slides.add_slide(BLANK)
chrome(s, "Hosted backend · Admin", "Admin security & console", accent=AMBER)
two_col(s, "Admin authentication", [
    "Tenant-admin API (/admin/*) behind a static token OR OIDC SSO",
    "OIDC: IdP-issued JWT, audience-checked, optional role mapping",
    "Provision tenants; issue / rotate / revoke per-tenant API keys",
    "Every admin mutation written to a separate admin-audit trail",
], "Graphical console", [
    "Opt-in static admin console served at /console",
    "Token-entry auth; holds no secrets (the API gate enforces)",
    "Browse tenants, keys, metrics & audit from a browser",
    "Off by default — enabled per deployment",
], accentL=AMBER, accentR=AMBER)

# 15 — section: flavors
s = prs.slides.add_slide(BLANK)
N[0] += 1
bg(s)
rect(s, Inches(0.9), Inches(3.05), Inches(2.2), Inches(0.06), EKG)
text(s, Inches(0.9), Inches(2.2), Inches(11), Inches(0.6), [[("THREE MARKET FLAVORS", 16, EKG, True, MONO)]])
text(s, Inches(0.88), Inches(2.7), Inches(11.5), Inches(1.2),
     [[("One core, productized three ways", 34, INK, True, SANS)]])
text(s, Inches(0.92), Inches(3.4), Inches(11.5), Inches(0.6),
     [[("Humanitarian / NGO   ·   Ontario EMS (regulated)   ·   Productized backend", 16, DIM, False, SANS)]])
footer(s, N[0])

# 16 — humanitarian
s = prs.slides.add_slide(BLANK)
chrome(s, "Flavor · Humanitarian / NGO", "Field documentation where the cloud isn't")
two_col(s, "Shipped on this line", [
    "Deployment context — device-wide operation tag + provenance banner",
    "Disaster/MCI mode — one toggle makes encryption mandatory + command roll-up",
    "Kiosk defaults — assign-operator prompt + 2-min auto-lock on shared devices",
    "Donor export — provenance-stamped CSV + date-range filter",
], "More", [
    "Retention presets — confirmed, step-up-gated purge window (off/30/90/180/365 d)",
    "Air-gapped packaging — one Docker stack (PWA + sync + Postgres) offline on a laptop",
    "Loadable language packs for any region; encrypted backups for low-infra handoff",
    "Intended use = documentation/coordination tool (light reg load)",
], accentL=EKG, accentR=EKG)

# 17 — ontario
s = prs.slides.add_slide(BLANK)
chrome(s, "Flavor · Ontario EMS (regulated)", "Toward a conformant ePCR", accent=BLUE)
two_col(s, "NEMSIS / OADS conformance", [
    "Section exporter maps captured data → NEMSIS v3.5 / OADS v4.0 shapes",
    "Deterministic offline XML serializer (shaped, gap-annotated)",
    "Pluggable validator (cardinality/datatype/value-set) + placeholder ruleset",
    "Capture panels: eResponse / eTimes / eCrew / eScene — gaps clear at runtime",
], "In-app & integration", [
    "Read-only Conformance view: live gaps + validator issues, in 4 languages",
    "Clearly labelled NOT certification (placeholder ruleset, rulesetSource shown)",
    "ONE ID / Ontario Health gateway scaffold (mTLS, server-side) for PCR $match",
    "Gated next: official-dictionary reconciliation + live ONE ID credentials",
], accentL=BLUE, accentR=BLUE)

# 18 — productized backend
s = prs.slides.add_slide(BLANK)
chrome(s, "Flavor · Productized backend", "The shared, hardened service line", accent=AMBER)
bullets(s, [
    "The multi-tenant sync service + admin security is the common spine both market branches reuse.",
    "Multi-tenant isolation, OIDC SSO + role-based admin, per-tenant rate limits & quota, incremental sync.",
    "EHR-access and admin audit trails — necessary for any regulated, multi-service deployment.",
    "Cascaded to every flavor branch, so a hardening fix lands everywhere.",
])

# 19 — status
s = prs.slides.add_slide(BLANK)
chrome(s, "Status", "Where things stand")
two_col(s, "Done & merged", [
    "Full offline PWA: capture, triage, vitals, handover, summaries",
    "Encryption vault · operators · step-up · tamper-evident audit",
    "Multi-tenant backend hardening + graphical admin console",
    "Humanitarian: all 5 roadmap items (context→MCI→export→retention→air-gap)",
    "Ontario: NEMSIS pipeline + capture panels + in-app conformance view",
], "Externally gated next", [
    "Ontario PR-2c — official NEMSIS/OADS dictionary + true XSD validation",
    "Ontario PR-4 — live ONE ID credentials + real mTLS client cert",
    "SOC 2 Type II / QMS / SaMD evidence (if intended use crosses the device line)",
    "CAD / hospital-EHR handover partnerships",
], accentL=EKG, accentR=AMBER)

# 20 — closing
s = prs.slides.add_slide(BLANK)
N[0] += 1
bg(s)
rect(s, 0, Inches(3.5), EMU_W, Inches(0.045), EKG)
text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0), [[("Document anywhere.", 40, INK, True, SANS)]])
text(s, Inches(0.92), Inches(3.7), Inches(11.5), Inches(0.6),
     [[("Offline-first by design · encrypted & audited · ready for the cloud only when you want it.", 17, EKG, False, SANS)]])
text(s, Inches(0.92), Inches(4.5), Inches(11.5), Inches(0.5),
     [[("TRIAGE-LINK — PWA · multi-tenant backend · Humanitarian & Ontario EMS flavors", 13, DIM, False, MONO)]])
footer(s, N[0])

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "TRIAGE-LINK-overview.pptx")
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
